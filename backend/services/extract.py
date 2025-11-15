import os
import pdfplumber
from pptx import Presentation

UPLOAD_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "data", "uploads"))
MAX_CHARS_HARD_LIMIT = 75_000  # ~25k tokens; safe for most 32k-context models

def _read_pdf_text(path: str) -> str:
    text = ""
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            try:
                page_text = page.extract_text() or ""
            except Exception:
                page_text = ""
            text += page_text + "\n"
    return text

def _read_pptx_text(path: str) -> str:
    text_chunks = []
    try:
        pres = Presentation(path)
        for slide in pres.slides:
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text") and shape.text:
                        text_chunks.append(shape.text)
                except Exception:
                    continue
    except Exception:
        return ""
    return "\n\n".join(text_chunks)

def _read_plain_text(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception:
        return ""

def read_document_text(filename: str, max_chars: int = MAX_CHARS_HARD_LIMIT) -> str:
    """
    Read text from an uploaded document (PDF, PPTX, or plain text).
    - By default, uses full document up to MAX_CHARS_HARD_LIMIT characters.
    - If max_chars is None → no trimming.
    - If max_chars is provided (e.g., 8000), that per-call limit is applied.
    """
    path = os.path.join(UPLOAD_DIR, filename)
    path = os.path.abspath(path)

    if not os.path.exists(path):
        return ""

    _, ext = os.path.splitext(path)
    ext = ext.lower()

    if ext == ".pdf":
        text = _read_pdf_text(path)
    elif ext in (".pptx", ".ppt"):
        text = _read_pptx_text(path)
    else:
        # fallback: treat as plain text
        text = _read_plain_text(path)

    if not text:
        return ""

    # Clean basic junk
    text = text.replace("\x00", "").strip()
    if not text:
        return ""

    # Recommended whitespace normalization
    text = "\n".join(line.strip() for line in text.splitlines() if line.strip())

    # If caller explicitly says "no limit", return full text
    if max_chars is None:
        return text
    
    # Enforce per-call OR hard limit
    if len(text) > max_chars:
        # Keep both beginning and end for context
        half = max_chars // 2
        text = text[:half] + "\n\n[... trimmed for length ...]\n\n" + text[-half:]

    return text
